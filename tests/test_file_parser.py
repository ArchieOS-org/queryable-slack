"""
Tests for file parsing functionality.

Tests CSV, XLSX, PPTX, ZIP, PDF, DOCX, TXT, and metadata extraction.
"""

import pytest
from pathlib import Path
import csv
import json
import zipfile
from conductor.file_parser import extract_text_from_file, extract_file_metadata


class TestFileMetadata:
    """Tests for extract_file_metadata function."""

    def test_extract_metadata_txt_file(self, tmp_path):
        """Test metadata extraction from a text file."""
        print("\n📄 Testing metadata extraction from text file...")
        test_file = tmp_path / "test.txt"
        test_file.write_text("Hello, World!")
        print(f"   Created test file: {test_file}")
        
        print("   Extracting metadata...")
        metadata = extract_file_metadata(test_file)
        print(f"   ✅ Metadata extracted: {metadata['filename']} ({metadata['file_type']}, {metadata['size']} bytes)")
        
        assert metadata["filename"] == "test.txt"
        assert metadata["file_type"] == "txt"
        assert metadata["size"] > 0
        assert metadata["extension"] == "txt"
        assert metadata["created"] is not None
        assert metadata["modified"] is not None
        print("   ✅ All assertions passed!")

    def test_extract_metadata_nonexistent_file(self):
        """Test that metadata extraction raises error for nonexistent file."""
        fake_path = Path("/nonexistent/file.txt")
        
        with pytest.raises(FileNotFoundError):
            extract_file_metadata(fake_path)

    def test_extract_metadata_infers_type_from_extension(self, tmp_path):
        """Test that file type is inferred from extension."""
        test_file = tmp_path / "document.pdf"
        test_file.write_text("fake pdf content")
        
        metadata = extract_file_metadata(test_file)
        
        assert metadata["file_type"] == "pdf"


class TestCSVParsing:
    """Tests for CSV file parsing."""

    def test_parse_simple_csv(self, tmp_path):
        """Test parsing a simple CSV file."""
        print("\n📊 Testing CSV file parsing...")
        csv_file = tmp_path / "test.csv"
        print(f"   Creating CSV file: {csv_file}")
        with open(csv_file, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Name", "Age", "City"])
            writer.writerow(["Alice", "30", "New York"])
            writer.writerow(["Bob", "25", "San Francisco"])
        print("   ✅ CSV file created with 3 rows")
        
        print("   Extracting text from CSV...")
        content = extract_text_from_file(csv_file)
        print(f"   ✅ Extracted {len(content)} characters")
        
        print("   Verifying content...")
        assert "Name" in content
        assert "Age" in content
        assert "Alice" in content
        assert "Bob" in content
        assert "30" in content
        print("   ✅ All CSV content assertions passed!")

    def test_parse_csv_with_special_characters(self, tmp_path):
        """Test parsing CSV with special characters."""
        csv_file = tmp_path / "test.csv"
        with open(csv_file, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Name", "Description"])
            writer.writerow(["Test", "Contains, commas and \"quotes\""])
        
        content = extract_text_from_file(csv_file)
        
        assert "Test" in content
        assert "commas" in content


class TestXLSXParsing:
    """Tests for XLSX file parsing."""

    def test_parse_xlsx_file(self, tmp_path):
        """Test parsing an XLSX file."""
        print("\n📈 Testing XLSX file parsing...")
        try:
            import openpyxl
            print("   ✅ openpyxl is available")
        except ImportError:
            print("   ⏭️  Skipping: openpyxl not installed")
            pytest.skip("openpyxl not installed")
        
        xlsx_file = tmp_path / "test.xlsx"
        print(f"   Creating XLSX file: {xlsx_file}")
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet["A1"] = "Name"
        sheet["B1"] = "Value"
        sheet["A2"] = "Test"
        sheet["B2"] = 42
        workbook.save(xlsx_file)
        print("   ✅ XLSX file created")
        
        print("   Extracting text from XLSX...")
        content = extract_text_from_file(xlsx_file)
        print(f"   ✅ Extracted {len(content)} characters")
        
        print("   Verifying content...")
        assert "Name" in content
        assert "Value" in content
        assert "Test" in content
        assert "42" in content
        print("   ✅ All XLSX content assertions passed!")


class TestPPTXParsing:
    """Tests for PPTX file parsing."""

    def test_parse_pptx_file(self, tmp_path):
        """Test parsing a PPTX file."""
        print("\n📽️  Testing PPTX file parsing...")
        try:
            from pptx import Presentation
            print("   ✅ python-pptx is available")
        except ImportError:
            print("   ⏭️  Skipping: python-pptx not installed")
            pytest.skip("python-pptx not installed")
        
        pptx_file = tmp_path / "test.pptx"
        print(f"   Creating PPTX file: {pptx_file}")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Slide"
        content = slide.placeholders[1]
        content.text = "This is test content"
        prs.save(pptx_file)
        print("   ✅ PPTX file created with 1 slide")
        
        print("   Extracting text from PPTX...")
        content = extract_text_from_file(pptx_file)
        print(f"   ✅ Extracted {len(content)} characters")
        
        print("   Verifying content...")
        assert "Test Slide" in content
        assert "test content" in content
        print("   ✅ All PPTX content assertions passed!")


class TestZIPParsing:
    """Tests for ZIP file parsing."""

    def test_parse_zip_file(self, tmp_path):
        """Test parsing a ZIP file with text contents."""
        print("\n📦 Testing ZIP file parsing...")
        zip_file = tmp_path / "test.zip"
        print(f"   Creating ZIP file: {zip_file}")
        
        # Create a ZIP with text files
        with zipfile.ZipFile(zip_file, "w") as zf:
            zf.writestr("file1.txt", "Content of file 1")
            zf.writestr("file2.csv", "Name,Age\nAlice,30")
            zf.writestr("binary.bin", b"\x00\x01\x02")
        print("   ✅ ZIP file created with 3 files")
        
        print("   Extracting text from ZIP...")
        content = extract_text_from_file(zip_file)
        print(f"   ✅ Extracted {len(content)} characters")
        
        print("   Verifying content...")
        assert "file1.txt" in content
        assert "file2.csv" in content
        assert "Content of file 1" in content
        assert "Alice" in content
        print("   ✅ All ZIP content assertions passed!")


class TestTxtParsing:
    """Tests for TXT file parsing."""

    def test_parse_utf8_txt(self, tmp_path):
        """Test parsing UTF-8 text file."""
        txt_file = tmp_path / "test.txt"
        txt_file.write_text("Hello, World! 🌍", encoding="utf-8")
        
        content = extract_text_from_file(txt_file)
        
        assert "Hello, World!" in content
        assert "🌍" in content

    def test_parse_latin1_txt(self, tmp_path):
        """Test parsing Latin-1 encoded text file."""
        txt_file = tmp_path / "test.txt"
        txt_file.write_bytes("Café résumé".encode("latin-1"))
        
        content = extract_text_from_file(txt_file)
        
        assert "Café" in content or "Caf" in content


class TestImageHandling:
    """Tests for image file handling (metadata only)."""

    def test_image_returns_metadata(self, tmp_path):
        """Test that images return structured metadata."""
        # Create a fake image file (just empty file with .png extension)
        img_file = tmp_path / "test.png"
        img_file.write_bytes(b"fake image data")
        
        content = extract_text_from_file(img_file)
        
        assert "[SKIPPED:" in content or "[SKIPPED:" in content
        assert "test.png" in content
        assert "png" in content.lower()


class TestVideoHandling:
    """Tests for video file handling."""

    def test_video_returns_metadata_or_processed(self, tmp_path):
        """Test that videos return metadata or processed content."""
        # Create a fake video file
        video_file = tmp_path / "test.mp4"
        video_file.write_bytes(b"fake video data")
        
        content = extract_text_from_file(video_file)
        
        # Should either be skipped with metadata or processed
        assert "[SKIPPED:" in content or "[VIDEO_PROCESSED:" in content
        assert "test.mp4" in content or "mp4" in content.lower()


class TestAudioHandling:
    """Tests for audio file handling."""

    def test_audio_returns_metadata_or_processed(self, tmp_path):
        """Test that audio files return metadata or processed content."""
        # Create a fake audio file
        audio_file = tmp_path / "test.mp3"
        audio_file.write_bytes(b"fake audio data")
        
        content = extract_text_from_file(audio_file)
        
        # Should either be skipped with metadata or processed
        assert "[SKIPPED:" in content or "[AUDIO_PROCESSED:" in content
        assert "test.mp3" in content or "mp3" in content.lower()


class TestErrorHandling:
    """Tests for error handling in file parsing."""

    def test_nonexistent_file_raises_error(self):
        """Test that nonexistent file raises FileNotFoundError."""
        fake_path = Path("/nonexistent/file.pdf")
        
        with pytest.raises(FileNotFoundError):
            extract_text_from_file(fake_path)

    def test_corrupt_file_returns_error_message(self, tmp_path):
        """Test that corrupt files return error message instead of crashing."""
        # Create a file that looks like PDF but isn't
        fake_pdf = tmp_path / "corrupt.pdf"
        fake_pdf.write_text("This is not a real PDF file")
        
        content = extract_text_from_file(fake_pdf)
        
        # Should return error message, not crash
        assert isinstance(content, str)
        assert len(content) > 0

