"""
Tests to verify file type queries work correctly.

Tests that queries for specific file types (CSV, PPTX, etc.) can find relevant documents.
"""

import pytest
from pathlib import Path
import chromadb
from conductor.ingest import enrich_session_with_files, store_sessions_in_chromadb
from conductor.models import Session, SlackMessage
from conductor.ask import query_chromadb
from datetime import datetime
import csv
import tempfile
import shutil


class TestFileTypeQueryMatching:
    """Tests that file type information is prominent enough for semantic search."""

    def test_csv_file_type_prominence(self, tmp_path):
        """Test that CSV files are labeled prominently for semantic search."""
        print("\n🔍 Testing CSV file type prominence in enriched transcript...")
        
        # Create test CSV
        attachments_dir = tmp_path / "attachments"
        attachments_dir.mkdir(parents=True)
        csv_file = attachments_dir / "F123-test.csv"
        
        with open(csv_file, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Property", "Price"])
            writer.writerow(["123 Main St", "$500,000"])
        print(f"   Created CSV file: {csv_file}")
        
        # Create session
        session = Session(
            session_id="test-csv-session",
            start_time=datetime.now(),
            end_time=datetime.now(),
            channel_name="test-channel",
            conversation_type="channel",
            transcript="User: Check this CSV file",
            enriched_transcript="",
            file_count=0,
            message_count=1,
        )
        
        messages = [
            SlackMessage(
                ts="1234567890.123456",
                user="U12345",
                text="Check this CSV file",
                type="message",
                files=[{"id": "F123", "name": "test.csv", "filetype": "csv"}]
            )
        ]
        
        print("   Enriching session...")
        enriched = enrich_session_with_files(session, messages, tmp_path)
        
        # Verify file type information is prominent
        print("   Verifying file type prominence...")
        assert "CSV" in enriched.enriched_transcript.upper(), "CSV should appear in uppercase"
        assert "CSV file spreadsheet" in enriched.enriched_transcript or "CSV FILE" in enriched.enriched_transcript, "CSV file type label should be present"
        assert "File type: CSV" in enriched.enriched_transcript, "File type metadata should include CSV"
        print("   ✅ CSV file type is prominently labeled!")

    def test_pptx_file_type_prominence(self, tmp_path):
        """Test that PPTX files are labeled prominently for semantic search."""
        print("\n🔍 Testing PPTX file type prominence...")
        
        try:
            from pptx import Presentation
        except ImportError:
            print("   ⏭️  Skipping: python-pptx not installed")
            pytest.skip("python-pptx not installed")
        
        # Create test PPTX
        attachments_dir = tmp_path / "attachments"
        attachments_dir.mkdir(parents=True)
        pptx_file = attachments_dir / "F456-presentation.pptx"
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Presentation"
        prs.save(pptx_file)
        print(f"   Created PPTX file: {pptx_file}")
        
        # Create session
        session = Session(
            session_id="test-pptx-session",
            start_time=datetime.now(),
            end_time=datetime.now(),
            channel_name="test-channel",
            conversation_type="channel",
            transcript="User: Check this presentation",
            enriched_transcript="",
            file_count=0,
            message_count=1,
        )
        
        messages = [
            SlackMessage(
                ts="1234567890.123456",
                user="U12345",
                text="Check this presentation",
                type="message",
                files=[{"id": "F456", "name": "presentation.pptx", "filetype": "pptx"}]
            )
        ]
        
        print("   Enriching session...")
        enriched = enrich_session_with_files(session, messages, tmp_path)
        
        # Verify file type information is prominent
        print("   Verifying file type prominence...")
        assert "PPTX" in enriched.enriched_transcript.upper(), "PPTX should appear in uppercase"
        assert "PowerPoint" in enriched.enriched_transcript or "PPTX FILE" in enriched.enriched_transcript, "PowerPoint/PPTX label should be present"
        assert "File type: PPTX" in enriched.enriched_transcript, "File type metadata should include PPTX"
        print("   ✅ PPTX file type is prominently labeled!")

    def test_semantic_search_finds_file_types(self, tmp_path):
        """Test that semantic search can find files by type."""
        print("\n🔍 Testing semantic search for file types...")
        
        # Create test database
        db_path = tmp_path / "test_db"
        db_path.mkdir()
        
        # Create CSV file
        attachments_dir = tmp_path / "test_channel" / "attachments"
        attachments_dir.mkdir(parents=True)
        csv_file = attachments_dir / "F789-data.csv"
        
        with open(csv_file, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Name", "Value"])
            writer.writerow(["Test", "42"])
        
        # Create session with CSV
        session = Session(
            session_id="csv-search-test",
            start_time=datetime.now(),
            end_time=datetime.now(),
            channel_name="test_channel",
            conversation_type="channel",
            transcript="User: Here's the data",
            enriched_transcript="",
            file_count=0,
            message_count=1,
        )
        
        messages = [
            SlackMessage(
                ts="1234567890.123456",
                user="U12345",
                text="Here's the data",
                type="message",
                files=[{"id": "F789", "name": "data.csv", "filetype": "csv"}]
            )
        ]
        
        print("   Enriching session with CSV...")
        enriched = enrich_session_with_files(session, messages, tmp_path / "test_channel")
        
        # Store in ChromaDB
        print("   Storing in ChromaDB...")
        store_sessions_in_chromadb([enriched], db_path=db_path)
        
        # Query for CSV files
        print("   Querying for CSV files...")
        results = query_chromadb("CSV files or spreadsheet data", db_path=db_path, n_results=5)
        
        # Verify results
        print("   Verifying search results...")
        assert results.get("documents") and len(results["documents"][0]) > 0, "Should find at least one result"
        
        # Check if CSV-related content is in results
        found_csv = False
        for doc in results["documents"][0]:
            if "CSV" in doc.upper() or "spreadsheet" in doc.lower() or "File type: CSV" in doc:
                found_csv = True
                break
        
        assert found_csv, "Query results should contain CSV-related content"
        print("   ✅ Semantic search successfully finds CSV files!")

