"""
File parsing wrapper for LangChain/Unstructured loaders.

Handles PDF, DOCX, TXT, CSV, XLSX, PPTX, ZIP file extraction with error handling.
Also extracts metadata for all file types.
"""

import logging
import os
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Dict, Optional, Any

import json

# Lazy imports for heavy dependencies - only load when needed
_pandas = None
_py_pdf_loader = None
_unstructured_loader = None
_pptx_presentation = None
_openpyxl = None
_xlrd = None

def _get_pandas():
    """Lazy load pandas."""
    global _pandas
    if _pandas is None:
        import pandas as pd
        _pandas = pd
    return _pandas

def _get_pypdf_loader():
    """Lazy load PyPDFLoader."""
    global _py_pdf_loader
    if _py_pdf_loader is None:
        from langchain_community.document_loaders import PyPDFLoader
        _py_pdf_loader = PyPDFLoader
    return _py_pdf_loader

def _get_unstructured_loader():
    """Lazy load UnstructuredFileLoader (with deprecation warning suppression)."""
    global _unstructured_loader
    if _unstructured_loader is None:
        import warnings
        # Suppress LangChain deprecation warnings for UnstructuredFileLoader
        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", category=DeprecationWarning, module="langchain")
            from langchain_community.document_loaders import UnstructuredFileLoader
            _unstructured_loader = UnstructuredFileLoader
    return _unstructured_loader

def _get_pptx():
    """Lazy load python-pptx."""
    global _pptx_presentation
    if _pptx_presentation is None:
        from pptx import Presentation
        _pptx_presentation = Presentation
    return _pptx_presentation

def _get_openpyxl():
    """Lazy load openpyxl."""
    global _openpyxl
    if _openpyxl is None:
        try:
            import openpyxl
            _openpyxl = openpyxl
        except ImportError:
            _openpyxl = None
    return _openpyxl

def _get_xlrd():
    """Lazy load xlrd."""
    global _xlrd
    if _xlrd is None:
        try:
            import xlrd
            _xlrd = xlrd
        except ImportError:
            _xlrd = None
    return _xlrd

# Import video, audio, and image processors (optional dependencies)
try:
    from conductor.video_processor import process_video_content
except ImportError:
    process_video_content = None

try:
    from conductor.audio_processor import process_audio_content
except ImportError:
    process_audio_content = None

try:
    from conductor.image_processor import process_image_content
except ImportError:
    process_image_content = None

logger = logging.getLogger(__name__)


def extract_file_metadata(file_path: Path, file_type: Optional[str] = None) -> Dict[str, Any]:
    """
    Extract metadata from a file.

    Args:
        file_path: Path to the file
        file_type: File type hint (e.g., "pdf", "docx", "txt"). If None, inferred from extension.

    Returns:
        Dictionary containing file metadata:
        - filename: str
        - file_type: str
        - size: int (bytes)
        - extension: str
        - created: Optional[datetime]
        - modified: Optional[datetime]
    """
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    # Infer file type from extension if not provided
    if file_type is None:
        file_type = file_path.suffix.lower().lstrip(".")

    file_type = file_type.lower()

    # Get file stats
    stat = file_path.stat()
    size = stat.st_size

    # Get timestamps
    created = datetime.fromtimestamp(stat.st_ctime) if hasattr(stat, 'st_ctime') else None
    modified = datetime.fromtimestamp(stat.st_mtime) if hasattr(stat, 'st_mtime') else None

    metadata = {
        "filename": file_path.name,
        "file_type": file_type,
        "size": size,
        "extension": file_path.suffix.lower().lstrip("."),
        "created": created.isoformat() if created else None,
        "modified": modified.isoformat() if modified else None,
    }

    return metadata


def extract_text_from_file(file_path: Path, file_type: Optional[str] = None) -> str:
    """
    Extract text content from a file using appropriate loader.

    Supports PDF, DOCX, TXT, CSV, XLSX, XLS, PPTX, and ZIP files.
    Images, videos, and audio files return structured metadata instead of content.
    Unsupported types return structured metadata.

    Args:
        file_path: Path to the file to extract text from
        file_type: File type hint (e.g., "pdf", "docx", "txt", "csv", "xlsx", "pptx", "zip").
                   If None, inferred from extension.

    Returns:
        Extracted text content, structured metadata for skipped files, or error message if extraction fails

    Raises:
        FileNotFoundError: If file doesn't exist
    """
    if not file_path.exists():
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    # Infer file type from extension if not provided
    if file_type is None:
        file_type = file_path.suffix.lower().lstrip(".")

    file_type = file_type.lower()

    # Handle PDF files
    if file_type == "pdf":
        try:
            PyPDFLoader = _get_pypdf_loader()
            loader = PyPDFLoader(str(file_path))
            docs = loader.load()
            text_content = "\n\n".join([doc.page_content for doc in docs])
            logger.debug(f"Extracted {len(docs)} pages from PDF: {file_path.name}")
            return text_content
        except Exception as e:
            logger.warning(f"Failed to parse PDF {file_path.name}: {e}")
            return f"[ERROR: Could not parse file {file_path.name}]"

    # Handle DOCX files
    if file_type == "docx":
        try:
            UnstructuredFileLoader = _get_unstructured_loader()
            loader = UnstructuredFileLoader(str(file_path))
            docs = loader.load()
            text_content = "\n\n".join([doc.page_content for doc in docs])
            logger.debug(f"Extracted content from DOCX: {file_path.name}")
            return text_content
        except Exception as e:
            logger.warning(f"Failed to parse DOCX {file_path.name}: {e}")
            return f"[ERROR: Could not parse file {file_path.name}]"

    # Handle TXT files
    if file_type == "txt":
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text_content = f.read()
            logger.debug(f"Read text file: {file_path.name}")
            return text_content
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, "r", encoding="latin-1") as f:
                    text_content = f.read()
                logger.debug(f"Read text file with latin-1 encoding: {file_path.name}")
                return text_content
            except Exception as e:
                logger.warning(f"Failed to read text file {file_path.name}: {e}")
                return f"[ERROR: Could not parse file {file_path.name}]"
        except Exception as e:
            logger.warning(f"Failed to read text file {file_path.name}: {e}")
            return f"[ERROR: Could not parse file {file_path.name}]"

    # Handle CSV files
    if file_type == "csv":
        try:
            # Try reading with different encodings and delimiters
            pd = _get_pandas()
            df = pd.read_csv(file_path, encoding="utf-8", on_bad_lines="skip")
            # Convert DataFrame to readable text format
            text_content = df.to_string(index=False)
            logger.debug(f"Extracted content from CSV: {file_path.name} ({len(df)} rows)")
            return text_content
        except UnicodeDecodeError:
            try:
                pd = _get_pandas()
                df = pd.read_csv(file_path, encoding="latin-1", on_bad_lines="skip")
                text_content = df.to_string(index=False)
                logger.debug(f"Extracted content from CSV with latin-1 encoding: {file_path.name}")
                return text_content
            except Exception as e:
                logger.warning(f"Failed to parse CSV {file_path.name}: {e}")
                return f"[ERROR: Could not parse file {file_path.name}]"
        except Exception as e:
            logger.warning(f"Failed to parse CSV {file_path.name}: {e}")
            return f"[ERROR: Could not parse file {file_path.name}]"

    # Handle XLSX files
    if file_type == "xlsx":
        openpyxl = _get_openpyxl()
        if openpyxl is None:
            logger.warning(f"openpyxl not installed, cannot parse XLSX: {file_path.name}")
            return f"[ERROR: openpyxl not installed, cannot parse file {file_path.name}]"
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():  # Only add non-empty rows
                        text_parts.append(row_text)
            text_content = "\n".join(text_parts)
            logger.debug(f"Extracted content from XLSX: {file_path.name} ({len(workbook.sheetnames)} sheets)")
            return text_content
        except Exception as e:
            logger.warning(f"Failed to parse XLSX {file_path.name}: {e}")
            return f"[ERROR: Could not parse file {file_path.name}]"

    # Handle XLS files (legacy Excel)
    if file_type == "xls":
        xlrd = _get_xlrd()
        if xlrd is None:
            logger.warning(f"xlrd not installed, cannot parse XLS: {file_path.name}")
            return f"[ERROR: xlrd not installed, cannot parse file {file_path.name}]"
        try:
            workbook = xlrd.open_workbook(file_path)
            text_parts = []
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                text_parts.append(f"=== Sheet: {sheet.name} ===")
                for row_idx in range(sheet.nrows):
                    row_values = sheet.row_values(row_idx)
                    row_text = "\t".join(str(cell) if cell else "" for cell in row_values)
                    if row_text.strip():  # Only add non-empty rows
                        text_parts.append(row_text)
            text_content = "\n".join(text_parts)
            logger.debug(f"Extracted content from XLS: {file_path.name} ({workbook.nsheets} sheets)")
            return text_content
        except Exception as e:
            logger.warning(f"Failed to parse XLS {file_path.name}: {e}")
            return f"[ERROR: Could not parse file {file_path.name}]"

    # Handle PPTX files
    if file_type == "pptx":
        try:
            Presentation = _get_pptx()
            prs = Presentation(file_path)
            text_parts = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                text_parts.append(f"=== Slide {slide_idx} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        text_frame = shape.text_frame
                        # Get all text from the text frame
                        slide_text = text_frame.text
                        if slide_text.strip():
                            text_parts.append(slide_text)
            text_content = "\n\n".join(text_parts)
            logger.debug(f"Extracted content from PPTX: {file_path.name} ({len(prs.slides)} slides)")
            return text_content
        except Exception as e:
            logger.warning(f"Failed to parse PPTX {file_path.name}: {e}")
            return f"[ERROR: Could not parse file {file_path.name}]"

    # Handle ZIP files
    if file_type == "zip":
        try:
            text_parts = []
            with zipfile.ZipFile(file_path, "r") as zip_ref:
                file_list = zip_ref.namelist()
                text_parts.append(f"=== ZIP Contents: {len(file_list)} files ===")
                for file_name in file_list:
                    text_parts.append(f"File: {file_name}")
                    # Try to extract text from supported files inside ZIP
                    try:
                        # Check if it's a supported text file type
                        inner_ext = Path(file_name).suffix.lower().lstrip(".")
                        if inner_ext in ["txt", "csv", "json", "xml"]:
                            with zip_ref.open(file_name) as inner_file:
                                inner_content = inner_file.read()
                                try:
                                    inner_text = inner_content.decode("utf-8")
                                    text_parts.append(f"  Content:\n{inner_text}")
                                except UnicodeDecodeError:
                                    text_parts.append(f"  [Binary or unsupported encoding]")
                    except Exception as e:
                        logger.debug(f"Could not extract content from {file_name} in ZIP: {e}")
                        text_parts.append(f"  [Could not extract content]")
            text_content = "\n\n".join(text_parts)
            logger.debug(f"Extracted contents from ZIP: {file_path.name}")
            return text_content
        except Exception as e:
            logger.warning(f"Failed to parse ZIP {file_path.name}: {e}")
            return f"[ERROR: Could not parse file {file_path.name}]"

    # Handle images - process with MLX-VLM for descriptions
    image_types = {"png", "jpg", "jpeg", "gif", "bmp", "webp", "svg"}
    if file_type in image_types:
        metadata = extract_file_metadata(file_path, file_type)
        
        if process_image_content is not None:
            try:
                image_content = process_image_content(file_path)
                description_result = image_content.get("description_result", {})
                
                # Check if processing was successful
                if description_result.get("success", False):
                    description = description_result.get("description", "")
                    return f"[IMAGE_PROCESSED: Description: {description} | Metadata: {json.dumps(metadata)}]"
                else:
                    # Processing failed but we have metadata - skip gracefully
                    error = description_result.get("error", "Unknown error")
                    logger.debug(f"Image processing skipped for {file_path.name}: {error}")
                    return f"[SKIPPED: Image file {metadata['filename']} (size: {metadata['size']} bytes, type: {metadata['file_type']}) - {error}]"
            except Exception as e:
                logger.warning(f"Failed to process image {file_path.name}: {e}")
                return f"[SKIPPED: Image file {metadata['filename']} (size: {metadata['size']} bytes, type: {metadata['file_type']}) - Error: {e}]"
        else:
            # Fallback to metadata-only if processor not available
            return f"[SKIPPED: Image file {metadata['filename']} (size: {metadata['size']} bytes, type: {metadata['file_type']}) - MLX-VLM not available]"

    # Handle videos - process with video processor if available
    video_types = {"mp4", "mov", "avi", "mkv", "webm", "flv", "wmv"}
    if file_type in video_types:
        metadata = extract_file_metadata(file_path, file_type)
        
        # Try to process video if processor is available
        if process_video_content is not None:
            try:
                logger.info(f"Processing video: {file_path.name}")
                video_result = process_video_content(
                    file_path,
                    extract_audio=True,
                    generate_description=True,
                    fps=1.0,  # Base: 1 frame per second (reduced for large files)
                    max_frames=30,  # Base: 30 frames max (reduced for large files)
                    max_duration_seconds=600.0,  # Process up to 10 minutes
                    max_file_size_mb=2000.0  # Allow up to 2GB with adaptive processing
                )
                
                # Build structured content from video processing
                content_parts = []
                content_parts.append(f"[VIDEO_PROCESSED: {metadata['filename']}]")
                content_parts.append(f"Metadata: {json.dumps(video_result.get('metadata', {}))}")
                
                if video_result.get('description'):
                    content_parts.append(f"Description: {video_result['description']}")
                
                if video_result.get('frame_count', 0) > 0:
                    content_parts.append(f"Frames extracted: {video_result['frame_count']}")
                
                if video_result.get('audio_path'):
                    content_parts.append(f"Audio extracted: {video_result['audio_path']}")
                    # Transcribe audio if available
                    if process_audio_content is not None:
                        try:
                            audio_result = process_audio_content(
                                Path(video_result['audio_path']),
                                max_duration_seconds=1800.0,  # Limit to 30 minutes
                                max_file_size_mb=200.0  # Limit to 200MB
                            )
                            if audio_result.get('success') and audio_result.get('text'):
                                content_parts.append(f"Audio transcription: {audio_result['text']}")
                        except Exception as e:
                            logger.debug(f"Could not transcribe video audio: {e}")
                
                return "\n".join(content_parts)
            except Exception as e:
                logger.warning(f"Video processing failed for {file_path.name}: {e}")
                # Fall back to metadata-only
                return f"[SKIPPED: Video file {metadata['filename']} (size: {metadata['size']} bytes, type: {metadata['file_type']}, processing failed: {e})]"
        else:
            # No video processor available, return metadata only
            logger.debug(f"Skipped video (no processor): {file_path.name}")
            return f"[SKIPPED: Video file {metadata['filename']} (size: {metadata['size']} bytes, type: {metadata['file_type']})]"

    # Handle audio - process with audio processor if available
    audio_types = {"mp3", "wav", "m4a", "aac", "ogg", "flac", "wma"}
    if file_type in audio_types:
        metadata = extract_file_metadata(file_path, file_type)
        
        # Try to transcribe audio if processor is available
        if process_audio_content is not None:
            try:
                logger.info(f"Processing audio: {file_path.name}")
                audio_result = process_audio_content(
                    file_path,
                    max_duration_seconds=1800.0,  # Limit to 30 minutes
                    max_file_size_mb=200.0  # Limit to 200MB
                )
                
                # Build structured content from audio processing
                content_parts = []
                content_parts.append(f"[AUDIO_PROCESSED: {metadata['filename']}]")
                content_parts.append(f"Metadata: {json.dumps(metadata)}")
                
                if audio_result.get('success') and audio_result.get('text'):
                    content_parts.append(f"Transcription: {audio_result['text']}")
                    if audio_result.get('language'):
                        content_parts.append(f"Language: {audio_result['language']}")
                    if audio_result.get('segments'):
                        content_parts.append(f"Segments: {len(audio_result['segments'])}")
                else:
                    content_parts.append(f"Transcription failed: {audio_result.get('error', 'Unknown error')}")
                
                return "\n".join(content_parts)
            except Exception as e:
                logger.warning(f"Audio processing failed for {file_path.name}: {e}")
                # Fall back to metadata-only
                return f"[SKIPPED: Audio file {metadata['filename']} (size: {metadata['size']} bytes, type: {metadata['file_type']}, processing failed: {e})]"
        else:
            # No audio processor available, return metadata only
            logger.debug(f"Skipped audio (no processor): {file_path.name}")
            return f"[SKIPPED: Audio file {metadata['filename']} (size: {metadata['size']} bytes, type: {metadata['file_type']})]"

    # For other unsupported types, return structured metadata
    logger.debug(f"Unsupported file type '{file_type}': {file_path.name}")
    metadata = extract_file_metadata(file_path, file_type)
    return f"[SKIPPED: Unsupported file type {metadata['file_type']} - {metadata['filename']} (size: {metadata['size']} bytes)]"
